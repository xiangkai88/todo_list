from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[1]
VENDOR = ROOT / ".vendor"
if str(VENDOR) not in sys.path:
    sys.path.insert(0, str(VENDOR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = ROOT / "out"
MEDIA_DIR = ROOT / "extracted_media"
OUTPUT = OUT_DIR / "21-向凯-优化版.pptx"

PRIMARY = RGBColor(25, 71, 129)
PRIMARY_DARK = RGBColor(16, 43, 79)
ACCENT = RGBColor(37, 153, 114)
TEXT = RGBColor(47, 58, 74)
MUTED = RGBColor(106, 119, 138)
BG = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(220, 226, 235)
WARN = RGBColor(232, 141, 53)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"


WORK_ITEMS = [
    {
        "title": "分公司营销日报",
        "subtitle": "形成日度经营信息快报，支持市场动态跟踪和快速响应。",
        "bullets": [
            "整合当日关键经营指标，提升信息传递时效性。",
            "为分公司营销决策提供统一数据口径和日常参照。",
            "沉淀日报模板，降低重复整理成本。",
        ],
        "image": "image12.png",
    },
    {
        "title": "日前报价方案",
        "subtitle": "围绕交易策略输出报价方案，增强投标准备的规范性。",
        "bullets": [
            "对市场价格区间和边界条件进行整理。",
            "支持业务侧在报价前快速完成方案校验。",
            "提升方案输出的完整度和可复用性。",
        ],
        "image": "image13.png",
    },
    {
        "title": "分公司现货交易简报",
        "subtitle": "用简报方式汇总交易结果，便于经营复盘和过程留痕。",
        "bullets": [
            "梳理交易执行情况与关键结果。",
            "突出重要异常、机会点及后续关注事项。",
            "形成可周转、可沉淀的现货信息材料。",
        ],
        "image": "image14.png",
    },
    {
        "title": "分公司日前市场复盘分析",
        "subtitle": "以复盘视角看价格与策略，帮助团队优化下一步动作。",
        "bullets": [
            "从市场表现、执行结果和偏差原因三个层面复盘。",
            "补充对重点时段和关键节点的解释说明。",
            "为后续策略迭代提供依据。",
        ],
        "image": "image15.png",
    },
    {
        "title": "现货内部对标表",
        "subtitle": "建立横向对比机制，提升内部经营分析的可视化程度。",
        "bullets": [
            "统一内部对标口径，方便不同主体间比较。",
            "识别优势项、短板项和改进空间。",
            "为阶段性评估和经验复制提供支撑。",
        ],
        "image": "image16.png",
    },
    {
        "title": "全省现货结算情况",
        "subtitle": "跟踪区域结算表现，辅助研判经营结果与市场影响。",
        "bullets": [
            "围绕结算结果进行集中整理和直观呈现。",
            "帮助业务快速识别变化趋势与结构特征。",
            "增强对区域市场结算逻辑的理解。",
        ],
        "image": "image17.png",
    },
    {
        "title": "超额获利回收",
        "subtitle": "针对政策及结算变化进行专题梳理，提升专题研究输出质量。",
        "bullets": [
            "结合制度背景解释超额获利回收机制。",
            "对影响路径和关注重点进行条理化呈现。",
            "便于后续持续跟踪与内部沟通。",
        ],
        "image": "image18.png",
    },
    {
        "title": "运营中心周例会汇报材料",
        "subtitle": "服务例会汇报场景，提高材料表达的清晰度和管理价值。",
        "bullets": [
            "围绕周度重点事项提炼核心内容。",
            "兼顾经营结果、问题反馈和下步计划。",
            "支撑跨部门信息同步与会议沟通。",
        ],
        "image": "image19.png",
    },
    {
        "title": "年度竞价汇总表",
        "subtitle": "沉淀年度竞价过程与结果，方便横纵向回溯分析。",
        "bullets": [
            "统一汇总竞价数据，减少散点式查询。",
            "支持从年度维度观察经营策略变化。",
            "为阶段总结和经验复用提供基础材料。",
        ],
        "image": "image20.png",
    },
    {
        "title": "利用小时对比表",
        "subtitle": "以指标对比方式展示运行成效，增强结果解读的直观性。",
        "bullets": [
            "聚焦关键运行指标，形成对比视角。",
            "辅助发现差异原因和改善方向。",
            "服务于经营分析和内部沟通场景。",
        ],
        "image": "image21.png",
    },
]


def add_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, title, subtitle=None):
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.45)
    ).fill.solid()
    band = slide.shapes[-1]
    band.fill.fore_color.rgb = PRIMARY
    band.line.color.rgb = PRIMARY

    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(6.3), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT_CN
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = PRIMARY_DARK

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(7.0), Inches(0.35))
        p = sub.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = FONT_CN
        r.font.size = Pt(10.5)
        r.font.color.rgb = MUTED

    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.42), Inches(0.7), Inches(0.06)
    ).fill.solid()
    accent = slide.shapes[-1]
    accent.fill.fore_color.rgb = ACCENT
    accent.line.color.rgb = ACCENT


def style_shape(shape, fill_rgb, line_rgb=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb or fill_rgb


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT_CN
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullet_list(slide, x, y, w, h, bullets, color=TEXT, size=14):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {text}"
        p.font.name = FONT_CN
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(8)
    return box


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = WHITE
    hero.line.color.rgb = WHITE

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(4.2), Inches(7.5)
    )
    style_shape(left, PRIMARY_DARK)

    deco = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.7), Inches(0.0), Inches(0.5), Inches(7.5)
    )
    style_shape(deco, ACCENT)

    add_textbox(slide, Inches(0.7), Inches(1.2), Inches(2.8), Inches(0.5), "2025", 28, True, WHITE)
    add_textbox(slide, Inches(4.8), Inches(1.45), Inches(6.8), Inches(1.0), "年度工作总结报告", 28, True, PRIMARY_DARK)
    add_textbox(slide, Inches(4.8), Inches(2.35), Inches(5.5), Inches(0.5), "向凯", 18, True, ACCENT)
    add_textbox(slide, Inches(4.8), Inches(2.78), Inches(6.0), Inches(0.45), "聚焦成果沉淀、经营分析与流程支撑", 12, False, MUTED)

    for idx, txt in enumerate(["工作回顾", "关键亮点", "后续思考"]):
        y = 4.1 + idx * 0.72
        chip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(y), Inches(1.45), Inches(0.42)
        )
        style_shape(chip, BG, BG)
        add_textbox(slide, Inches(5.05), Inches(y + 0.07), Inches(2.2), Inches(0.25), txt, 11, False, PRIMARY_DARK)

    add_textbox(slide, Inches(4.8), Inches(6.5), Inches(3.6), Inches(0.3), "汇报日期 | 2026", 10.5, False, MUTED)
    return slide


def add_contents(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "目录", "CONTENTS")

    items = [
        ("01", "工作总结", "重点成果、材料沉淀与经营支撑"),
        ("02", "年度亮点", "能力提升、协同价值与工作方法"),
        ("03", "一点想法", "稳定基本盘、增强归属感、形成合力"),
    ]
    for idx, (num, title, desc) in enumerate(items):
        x = Inches(0.8 + idx * 4.2)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(3.4), Inches(2.6)
        )
        style_shape(card, WHITE, LINE)
        add_textbox(slide, x + Inches(0.25), Inches(2.28), Inches(0.9), Inches(0.4), num, 22, True, ACCENT)
        add_textbox(slide, x + Inches(0.25), Inches(2.8), Inches(2.4), Inches(0.4), title, 20, True, PRIMARY_DARK)
        add_textbox(slide, x + Inches(0.25), Inches(3.38), Inches(2.8), Inches(0.9), desc, 11.5, False, MUTED)
    return slide


def add_divider(prs, part_no, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY_DARK)
    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.32)
    )
    style_shape(stripe, ACCENT)
    add_textbox(slide, Inches(0.9), Inches(2.0), Inches(1.2), Inches(0.6), f"Part {part_no}", 18, True, WHITE)
    add_textbox(slide, Inches(0.9), Inches(2.75), Inches(4.8), Inches(0.8), title, 30, True, WHITE)
    add_textbox(slide, Inches(0.95), Inches(3.6), Inches(5.6), Inches(0.45), subtitle, 12, False, RGBColor(210, 220, 235))
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(1.5), Inches(3.8), Inches(3.8)
    )
    style_shape(box, PRIMARY, PRIMARY)
    add_textbox(slide, Inches(8.7), Inches(2.4), Inches(2.7), Inches(1.2), title, 24, True, WHITE, PP_ALIGN.CENTER)
    return slide


def add_work_slide(prs, item, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, item["title"], f"工作总结 | {index:02d}")

    left_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.75), Inches(3.2), Inches(4.95)
    )
    style_shape(left_panel, WHITE, LINE)
    add_textbox(slide, Inches(0.95), Inches(2.0), Inches(2.6), Inches(0.8), item["subtitle"], 17, True, PRIMARY_DARK)
    add_textbox(slide, Inches(0.95), Inches(3.0), Inches(2.4), Inches(0.28), "工作价值", 11, True, ACCENT)
    add_bullet_list(slide, Inches(0.95), Inches(3.35), Inches(2.55), Inches(2.7), item["bullets"], size=12.5)

    img_frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.15), Inches(1.75), Inches(8.45), Inches(4.95)
    )
    style_shape(img_frame, WHITE, LINE)
    add_textbox(slide, Inches(4.45), Inches(1.95), Inches(1.7), Inches(0.25), "成果示意", 11, True, MUTED)
    img_path = MEDIA_DIR / item["image"]
    slide.shapes.add_picture(str(img_path), Inches(4.45), Inches(2.25), width=Inches(7.8), height=Inches(4.05))

    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(6.95), Inches(11.95), Inches(0.18)
    )
    style_shape(footer, LINE)
    return slide


def add_highlights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "年度亮点总结", "从成果数量到工作方法，逐步形成稳定输出能力。")

    metrics = [
        ("10+", "核心材料沉淀", "覆盖日报、简报、复盘、结算、竞价等多个场景"),
        ("3类", "主要工作角色", "数据整理、经营分析、会议支撑"),
        ("持续化", "输出方式优化", "由单点任务转向模板化、复用化和结构化表达"),
    ]
    for idx, (num, title, desc) in enumerate(metrics):
        x = Inches(0.8 + idx * 4.15)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(3.55), Inches(2.2)
        )
        style_shape(card, WHITE, LINE)
        add_textbox(slide, x + Inches(0.28), Inches(2.22), Inches(1.6), Inches(0.45), num, 24, True, ACCENT)
        add_textbox(slide, x + Inches(0.28), Inches(2.78), Inches(2.5), Inches(0.4), title, 17, True, PRIMARY_DARK)
        add_textbox(slide, x + Inches(0.28), Inches(3.32), Inches(2.85), Inches(0.9), desc, 11.2, False, MUTED)

    add_textbox(slide, Inches(0.85), Inches(5.0), Inches(2.2), Inches(0.28), "阶段性收获", 11, True, ACCENT)
    add_bullet_list(
        slide,
        Inches(0.85),
        Inches(5.35),
        Inches(11.5),
        Inches(1.35),
        [
            "在频繁的事务性工作中保持稳定交付，材料质量和响应效率同步提升。",
            "逐步形成以模板、台账、汇总表为核心的工作沉淀方式，减少重复劳动。",
            "围绕市场、结算、竞价等主题持续积累经验，为后续分析和决策支撑打下基础。",
        ],
        size=12.5,
    )
    return slide


def add_ideas_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "一点想法", "稳定基本盘，提升归属感，让分散的个体形成更强的合力。")

    ideas = [
        (
            "1. 稳定基本盘",
            "公司已经积累了项目经验、业务方向和团队骨干，这些都是支撑持续发展的基本盘。只有基本盘稳，才能让业务拓展和制度优化更从容。",
            ACCENT,
        ),
        (
            "2. 理解位置差异",
            "不同岗位看到的信息和关注点不同，容易形成局部视角。加强沟通、共享背景，有助于减少误解，让判断更接近整体最优。",
            WARN,
        ),
        (
            "3. 增强归属感",
            "行业出差多、人员分散，归属感更需要被主动建设。让员工感受到被看见、被连接、被支持，团队的目标感才会更稳固。",
            PRIMARY,
        ),
    ]
    for idx, (title, body, color) in enumerate(ideas):
        y = Inches(1.9 + idx * 1.55)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.8), Inches(1.18)
        )
        style_shape(card, WHITE, LINE)
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), y + Inches(0.23), Inches(1.65), Inches(0.42)
        )
        style_shape(pill, color, color)
        add_textbox(slide, Inches(1.18), y + Inches(0.29), Inches(1.4), Inches(0.2), title, 10.5, True, WHITE)
        add_textbox(slide, Inches(2.95), y + Inches(0.2), Inches(8.9), Inches(0.78), body, 12.2, False, TEXT)
    return slide


def add_actions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "建议方向", "从组织协同、机制建设和人员体验三个层面增强团队凝聚力。")

    actions = [
        ("沟通更前置", "重点事项尽量形成固定同步节奏，让信息不依赖单点传递。"),
        ("沉淀更系统", "把日报、简报、复盘、台账继续标准化，减少重复投入。"),
        ("关怀更具体", "对长期外派、频繁出差人员增加连接机制和反馈渠道。"),
        ("目标更一致", "让个人工作与公司阶段目标建立更明确的对应关系。"),
    ]
    for idx, (title, desc) in enumerate(actions):
        col = idx % 2
        row = idx // 2
        x = Inches(0.9 + col * 6.0)
        y = Inches(2.0 + row * 2.0)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(5.45), Inches(1.45)
        )
        style_shape(card, WHITE, LINE)
        num = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.28), y + Inches(0.32), Inches(0.62), Inches(0.62)
        )
        style_shape(num, PRIMARY, PRIMARY)
        add_textbox(slide, x + Inches(0.49), y + Inches(0.47), Inches(0.2), Inches(0.2), str(idx + 1), 11, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(1.05), y + Inches(0.25), Inches(2.0), Inches(0.25), title, 14.5, True, PRIMARY_DARK)
        add_textbox(slide, x + Inches(1.05), y + Inches(0.67), Inches(3.95), Inches(0.5), desc, 11.5, False, MUTED)
    return slide


def add_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY_DARK)
    add_textbox(slide, Inches(4.0), Inches(2.45), Inches(5.4), Inches(0.8), "感谢聆听", 30, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.55), Inches(3.25), Inches(6.3), Inches(0.35), "Thank You", 14, False, RGBColor(201, 214, 231), PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.15), Inches(5.8), Inches(3.0), Inches(0.25), "汇报人：向凯", 11.5, False, WHITE, PP_ALIGN.CENTER)
    return slide


def build():
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "2025年度工作总结报告"
    prs.core_properties.author = "Codex"
    prs.core_properties.subject = "优化版汇报材料"

    add_cover(prs)
    add_contents(prs)
    add_divider(prs, "1", "工作总结", "围绕核心材料与经营支撑工作，梳理年度输出成果。")
    for idx, item in enumerate(WORK_ITEMS, start=1):
        add_work_slide(prs, item, idx)
    add_highlights(prs)
    add_divider(prs, "2", "一点想法", "围绕基本盘、组织协同与归属感，提出一些工作感受。")
    add_ideas_summary(prs)
    add_actions(prs)
    add_closing(prs)

    prs.save(str(OUTPUT))
    print(OUTPUT.relative_to(ROOT))


if __name__ == "__main__":
    build()
